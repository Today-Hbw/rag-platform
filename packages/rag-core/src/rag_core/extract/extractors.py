"""附件文本提取分派与各格式实现。忠实移植 download.py:extract_text_from_file 及 _extract_*。

外部工具路径经 ToolPaths（可配置/PATH 发现）注入；重依赖缺失则该分支返回 None。
"""

from __future__ import annotations

import logging
import os
import re
import subprocess
import zipfile
from tempfile import TemporaryDirectory

from ..settings import ExtractSettings, Settings, get_settings
from .garbled import is_text_garbled
from .tools import ToolPaths

logger = logging.getLogger(__name__)

__all__ = ["extract_text_from_file"]


# ==================== PDF ====================

def _pdf_via_pdftotext(filepath: str, tools: ToolPaths) -> str | None:
    try:
        result = subprocess.run(
            [tools.pdftotext_cmd, "-layout", filepath, "-"],
            capture_output=True,
            timeout=60,
        )
        if result.returncode == 0:
            text = result.stdout.decode("utf-8", errors="replace").strip()
            return text if len(text) >= 10 else None
    except FileNotFoundError:
        pass  # pdftotext 未安装
    except Exception as e:
        logger.warning("pdftotext fallback failed for %s: %s", filepath, e)
    return None


def _pdf_via_ocr(filepath: str, tools: ToolPaths, ex: ExtractSettings) -> str | None:
    try:
        import pytesseract
        from pdf2image import convert_from_path
    except ImportError:
        return None
    try:
        if tools.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tools.tesseract_cmd
        images = convert_from_path(
            filepath,
            first_page=1,
            last_page=ex.ocr_max_pages,
            dpi=ex.ocr_dpi,
            poppler_path=tools.poppler_bin,
        )
        texts = []
        for img in images:
            text = pytesseract.image_to_string(img, lang=ex.ocr_lang)
            if text and text.strip():
                texts.append(text.strip())
        result = "\n".join(texts).strip()
        if result:
            logger.info("PDF OCR extracted %d chars from %s", len(result), filepath)
        return result or None
    except Exception as e:
        logger.warning("PDF OCR failed for %s: %s", filepath, e)
        return None


def _extract_pdf(filepath: str, tools: ToolPaths, ex: ExtractSettings) -> str | None:
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber not installed, skipping PDF text extraction")
        return None
    texts = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
                for table in page.extract_tables():
                    for row in table:
                        texts.append("\t".join(str(c or "") for c in row))
    except Exception as e:
        logger.warning("PDF text extraction failed for %s: %s", filepath, e)
        return None
    result = "\n".join(texts).strip()
    if result:
        return result
    # pdfplumber 空（可能扫描件）→ pdftotext → OCR 兜底
    result = _pdf_via_pdftotext(filepath, tools)
    if result:
        return result
    result = _pdf_via_ocr(filepath, tools, ex)
    if result:
        return result
    logger.warning("PDF text extraction returned empty (likely scanned): %s", filepath)
    return None


# ==================== .doc（旧版 Word）====================

def _doc_via_com(filepath: str) -> str | None:
    """Windows COM 驱动 Word 提取；非 Windows/未装 pywin32 返回 None。"""
    if os.name != "nt":
        return None
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return None

    abs_path = os.path.abspath(filepath)
    word = None
    doc = None
    try:
        pythoncom.CoInitialize()
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = False
        doc = word.Documents.Open(abs_path, ReadOnly=True)
        text = doc.Content.Text
        doc.Close(False)
        doc = None
        word.Quit()
        word = None
        pythoncom.CoUninitialize()
        if text:
            text = text.replace(chr(13), chr(10)).replace(chr(7), chr(9)).replace(chr(12), chr(10))
            text = text.strip()
        return text if text and len(text) >= 10 else None
    except Exception as e:
        logger.warning("COM extraction failed for %s: %s", filepath, e)
    finally:
        try:
            if doc is not None:
                doc.Close(False)
        except Exception:
            pass
        try:
            if word is not None:
                word.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
    return None


def _doc_via_libreoffice(filepath: str, tools: ToolPaths) -> str | None:
    lo_bin = tools.libreoffice_bin
    if not lo_bin:
        return None
    try:
        with TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                [lo_bin, "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, filepath],
                capture_output=True,
                timeout=60,
            )
            if result.returncode != 0:
                return None
            txt_name = os.path.splitext(os.path.basename(filepath))[0] + ".txt"
            txt_path = os.path.join(tmpdir, txt_name)
            if not os.path.exists(txt_path):
                return None
            with open(txt_path, encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()
            return text if len(text) >= 10 else None
    except Exception as e:
        logger.warning("LibreOffice extraction failed for %s: %s", filepath, e)
    return None


def _doc_via_ole(filepath: str) -> str | None:
    try:
        import olefile
    except ImportError:
        return None
    try:
        if not olefile.isOleFile(filepath):
            return None
        ole = olefile.OleFileIO(filepath)
        data = ole.openstream("WordDocument").read()
        ole.close()
    except Exception:
        return None
    utf16_len = len(data) - (len(data) % 2)
    decoded = data[:utf16_len].decode("utf-16-le", errors="ignore")
    cleaned = re.sub(r"[^\t\n\r -~一-鿿　-〿＀-￯]", " ", decoded)
    cleaned = re.sub(r" {3,}", "  ", cleaned).strip()
    return cleaned if len(cleaned) >= 20 else None


def _extract_doc(filepath: str, tools: ToolPaths) -> str | None:
    # 1. Windows COM
    text = _doc_via_com(filepath)
    if text and not is_text_garbled(text):
        return text
    # 2. antiword（UTF-8，失败再 GBK 重解）
    try:
        result = subprocess.run(
            [tools.antiword_cmd, "-m", "UTF-8.txt", filepath],
            capture_output=True,
            timeout=30,
        )
        if result.returncode == 0:
            text = result.stdout.decode("utf-8", errors="replace").strip()
            if text and not is_text_garbled(text):
                return text
            if result.stdout and len(result.stdout) > 20:
                raw_text = result.stdout.decode("gbk", errors="replace").strip()
                if raw_text and not is_text_garbled(raw_text):
                    return raw_text
    except FileNotFoundError:
        pass  # antiword 未安装
    except Exception as e:
        logger.warning("antiword failed for %s: %s", filepath, e)
    # 3. LibreOffice
    text = _doc_via_libreoffice(filepath, tools)
    if text and not is_text_garbled(text):
        return text
    # 4. OLE 兜底
    return _doc_via_ole(filepath)


# ==================== Office / 其他 ====================

def _extract_docx(filepath: str) -> str | None:
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed, skipping DOCX")
        return None
    try:
        with open(filepath, "rb") as f:  # 二进制对象避 Windows 非 ASCII 路径 bug
            doc = Document(f)
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    texts.append(row_text)
        return "\n".join(texts) or None
    except Exception as e:
        logger.warning("DOCX text extraction failed for %s: %s", filepath, e)
        return None


def _extract_xlsx(filepath: str) -> str | None:
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl not installed, skipping XLSX")
        return None
    try:
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            texts.append(f"== 工作表: {sheet_name} ==")
            for row in ws.iter_rows(values_only=True):
                texts.append("\t".join(str(c or "") for c in row))
        wb.close()
    except Exception as e:
        logger.warning("XLSX text extraction failed for %s: %s", filepath, e)
        return None
    return "\n".join(texts)


def _extract_xls(filepath: str) -> str | None:
    try:
        import xlrd
    except ImportError:
        logger.warning("xlrd not installed, skipping XLS")
        return None
    try:
        wb = xlrd.open_workbook(filepath)
        texts = []
        for sheet in wb.sheets():
            texts.append(f"== 工作表: {sheet.name} ==")
            for row_idx in range(sheet.nrows):
                texts.append("\t".join(str(c or "") for c in sheet.row_values(row_idx)))
    except Exception as e:
        logger.warning("XLS text extraction failed for %s: %s", filepath, e)
        return None
    return "\n".join(texts)


def _extract_pptx(filepath: str) -> str | None:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, skipping PPTX")
        return None
    try:
        with open(filepath, "rb") as f:
            prs = Presentation(f)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
    except Exception as e:
        logger.warning("PPTX text extraction failed for %s: %s", filepath, e)
        return None
    return "\n".join(texts)


def _extract_ppt(filepath: str) -> str | None:
    try:
        import olefile
    except ImportError:
        logger.warning("olefile not installed, skipping PPT")
        return None
    try:
        if not olefile.isOleFile(filepath):
            logger.warning("Not a valid OLE2 file: %s", filepath)
            return None
        ole = olefile.OleFileIO(filepath)
        data = ole.openstream("PowerPoint Document").read()
        ole.close()
        utf16_len = len(data) - (len(data) % 2)
        decoded = data[:utf16_len].decode("utf-16-le", errors="ignore")
        cleaned = re.sub(r"[^ -~一-鿿　-〿＀-￯\n\r\t]", " ", decoded)
        cleaned = re.sub(r" {3,}", "  ", cleaned).strip()
        return cleaned if len(cleaned) >= 20 else None
    except Exception as e:
        logger.warning("PPT text extraction failed for %s: %s", filepath, e)
        return None


def _extract_txt(filepath: str) -> str | None:
    try:
        with open(filepath, encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.warning("Text read failed for %s: %s", filepath, e)
        return None


def _extract_zip(
    filepath: str, settings: Settings, tools: ToolPaths, ex: ExtractSettings
) -> str | None:
    try:
        if not zipfile.is_zipfile(filepath):
            logger.warning("Not a valid zip file: %s", filepath)
            return None
        texts = []
        with zipfile.ZipFile(filepath, "r") as zf:
            members = [
                m
                for m in zf.infolist()
                if not m.is_dir()
                and not os.path.basename(m.filename).startswith(".")
                and m.file_size < ex.zip_max_file_size_mb * 1024 * 1024
            ][: ex.zip_max_files]
            with TemporaryDirectory() as tmpdir:
                zf.extractall(tmpdir, members)
                for m in members:
                    inner_path = os.path.join(tmpdir, m.filename)
                    if not os.path.isfile(inner_path):
                        continue
                    inner_ext = os.path.splitext(m.filename)[1].lower().lstrip(".")
                    if inner_ext:
                        inner_text = extract_text_from_file(
                            inner_path, inner_ext, settings=settings, tools=tools
                        )
                        if inner_text:
                            texts.append(f"[{os.path.basename(m.filename)}]\n{inner_text}")
        return "\n\n".join(texts) or None
    except Exception as e:
        logger.warning("ZIP text extraction failed for %s: %s", filepath, e)
        return None


def extract_text_from_file(
    filepath: str,
    file_type: str,
    *,
    settings: Settings | None = None,
    tools: ToolPaths | None = None,
) -> str | None:
    """从附件提取纯文本。支持 pdf/xlsx/xls/docx/pptx/doc/ppt/txt/csv/zip。"""
    settings = settings or get_settings()
    ex = settings.extract
    tools = tools or ToolPaths.resolve(settings)
    ft = file_type.lower().lstrip(".")

    if ft == "pdf":
        return _extract_pdf(filepath, tools, ex)
    if ft == "xlsx":
        return _extract_xlsx(filepath)
    if ft == "xls":
        return _extract_xls(filepath)
    if ft == "docx":
        return _extract_docx(filepath)
    if ft == "pptx":
        return _extract_pptx(filepath)
    if ft in ("txt", "csv"):
        return _extract_txt(filepath)
    if ft == "doc":
        return _extract_doc(filepath, tools)
    if ft == "ppt":
        return _extract_ppt(filepath)
    if ft == "zip":
        return _extract_zip(filepath, settings, tools, ex)

    logger.warning(
        "Unsupported attachment type for text extraction: '%s' (filepath=%s)", ft, filepath
    )
    return None
